import copy
import os
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt, Cm


class PPTReportConfig:
    """보고서 생성에 필요한 상수, 스타일, 규칙 등을 관리하는 설정 클래스"""

    # 레이아웃 & 폰트 설정
    MARGIN_LEFT = Cm(1.0) # Default (사용처 확인 후 필요시 제거)
    MARGIN_LEFT_L1 = Cm(0.62)       # Level 1: Main Sections, Failed gene
    MARGIN_LEFT_COMMENT = Cm(0.52)  # Level 1.5: Comments
    MARGIN_LEFT_L2 = Cm(0.72)       # Level 2: Sub-sections
    MARGIN_LEFT_L3 = Cm(0.92)       # Level 3: Tables
    
    DEFAULT_WIDTH = Cm(19.05)
    BODY_TOP_START = Cm(4.5)
    BODY_BOTTOM_LIMIT = Cm(24.5) # [Changed] 기존 18.0cm에서 확장하여 하단 여백 정상화

    FONT_NAME = "Malgun Gothic"
    FONT_SIZE_TITLE = Pt(12)
    FONT_SIZE_HEADER = Pt(12)
    FONT_SIZE_BODY = Pt(9)

    COLOR_BLACK = RGBColor(0, 0, 0)
    COLOR_RED = RGBColor(255, 0, 0)
    COLOR_NAVY = RGBColor(0, 51, 102)
    COLOR_GRAY_BG = RGBColor(230, 230, 230)

    # 간격 설정
    SPACE_SECTION = Cm(0.1)  # 섹션 간 간격
    SPACE_TABLE_BOTTOM = Cm(0.1)  # 테이블 하단 간격
    SPACE_TITLE_BOTTOM = Cm(0.05)  # 제목 하단 간격

    # 스타일 프리셋 정의
    STYLES = {
        "clinical": {"color": COLOR_RED, "bold": True},
        "unknown": {"color": COLOR_BLACK, "bold": False}
    }

    MARKER_BIOMARKERS = "Other Biomarkers"

    # 섹션 시작점을 찾기 위한 텍스트 매핑
    SECTION_START_MARKERS = {
        "clinical": "1. Variants of clinical significance",
        "unknown": "2. Variants of unknown significance"
    }

    # 고지문 시작 마커
    MARKER_DISCLAIMER = "*본 기관의 유전자 정보 검색"

    # 데이터 매핑
    CLINICAL_INFO_MAPPING = {
        "검체 정보": "검체 정보", "성별": "성별", "나이": "나이",
        "Unit NO.": "Unit NO.", "환자명": "환자명", "채취 장기": "채취 장기",
        "원발 장기": "원발 장기", "진단": "진단", "의뢰의": "의뢰의",
        "의뢰의 소속": "의뢰의 소속", "검체 유형": "검체 유형",
        "검체의 적절성여부": "검체의 적절성여부", "검체접수일": "검체접수일",
        "결과보고일": "결과보고일",
        "Tumor Mutation Burden": "tmb",
        "Microsatellite Instability": "msi"
    }

    # 섹션 순서 및 키
    VARIANT_SECTIONS = [
        # 1. Clinical Significance
        {"key": "snv_clinical", "title": "SNVs & Indels", "type": "clinical", "prototype_key": "snv_clinical"},
        {"key": "fusion_clinical", "title": "Fusion gene", "type": "clinical", "prototype_key": "fusion_clinical"},
        {"key": "cnv_clinical", "title": "Copy number variation", "type": "clinical", "prototype_key": "cnv_clinical"},
        {"key": "lr_brca_clinical", "title": "Large rearrangements in BRCA1/2", "type": "clinical",
         "prototype_key": "lr_brca_clinical"},
        {"key": "splice_clinical", "title": "Splice variant", "type": "clinical", "prototype_key": "splice_clinical"},

        # 2. Unknown Significance
        {"key": "snv_unknown", "title": "SNVs & Indels", "type": "unknown", "prototype_key": "snv_clinical"},
        {"key": "fusion_unknown", "title": "Fusion gene", "type": "unknown", "prototype_key": "fusion_clinical"},
        {"key": "cnv_unknown", "title": "Copy number variation", "type": "unknown", "prototype_key": "cnv_clinical"},
        {"key": "lr_brca_unknown", "title": "Large rearrangements in BRCA1/2", "type": "unknown",
         "prototype_key": "lr_brca_clinical"},
        {"key": "splice_unknown", "title": "Splice variant", "type": "unknown", "prototype_key": "splice_clinical"}
    ]

    # 테이블 식별 규칙 (확장성 핵심)
    # required: 반드시 포함되어야 할 헤더 키워드 (OR 조건은 튜플로 묶음)
    # excluded: 포함되면 안 되는 헤더 키워드
    TABLE_IDENTIFICATION_RULES = [
        {
            "key": "snv_clinical",
            "required": ["GENE", ("MUTATION", "AA CHANGE")],
            "excluded": ["BURDEN"]  # TMB 제외
        },
        {
            "key": "fusion_clinical",
            "required": ["FUSION", "BREAKPOINT"],
            "excluded": []
        },
        {
            "key": "cnv_clinical",
            "required": ["COPY", "FOLD"],
            "excluded": ["EXON"]  # Large Rearrangement와 구분
        },
        {
            "key": "lr_brca_clinical",
            "required": [("BRCA", "EXON"), "FOLD"],  # BRCA 혹은 (EXON and FOLD) - 로직에서 처리
            "excluded": []
        },
        {
            "key": "splice_clinical",
            "required": [("SPLICE", "EXON"), "BREAKPOINT"],
            "excluded": []
        }
    ]

    # 삭제할 소제목 키워드
    REMOVE_KEYWORDS = ["SNVs", "Fusion", "Copy", "Splice", "Rearrangement", "Failed gene"]


class LayoutAnalyzer:
    def __init__(self, prs):
        self.prs = prs
        self.page_width = prs.slide_width
        self.page_height = prs.slide_height
        self.config = PPTReportConfig()

        # 위치 초기화
        self.body_top = self.config.BODY_TOP_START
        self.body_bottom = self.config.BODY_BOTTOM_LIMIT

        # 슬라이드별 하단 제한선 저장소 {slide_index: bottom_limit}
        self.slide_bottom_limits = {}

        self.existing_elements = {
            "prototypes": {},
            "disclaimer": None  # 고지문 요소 저장용
        }

        # 섹션별 시작 위치 저장소
        self.section_locations = {}
        self.slide_bottom_limits = {}
        self._analyze_all_slides()

    def _analyze_all_slides(self):
        for slide_idx, slide in enumerate(self.prs.slides):
            self._analyze_single_slide(slide, slide_idx)

    def _analyze_single_slide(self, slide, slide_idx):
        shapes_to_remove = []
        found_section_on_this_slide = False

        # 이 슬라이드에서 사용할 하단 한계선 초기화 (기본값: 전역 설정 18cm)
        current_slide_limit = self.config.BODY_BOTTOM_LIMIT

        # 1. 텍스트 분석 (섹션 마커 및 Footer)
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()

                # "Other Biomarker" 텍스트 감지 -> 하단 한계선 위로 올림
                if "- Other Biomarker" in text:
                    # 해당 텍스트의 상단(Top)에서 0.5cm 위를 한계선으로 설정
                    limit = shape.top - Cm(0.5)
                    # 기존 한계선보다 더 위쪽에 있으면(작은 값이면) 업데이트하여 안전 확보
                    if limit < current_slide_limit:
                        current_slide_limit = limit

                # 섹션 시작 마커 찾기
                for section_type, marker in self.config.SECTION_START_MARKERS.items():
                    if marker in text:
                        self.section_locations[section_type] = {
                            "slide_index": slide_idx,
                            "slide_id": slide.slide_id,
                            "top": shape.top + shape.height + self.config.SPACE_TITLE_BOTTOM,
                            "has_title": True,
                            "title_shape": shape  # 제거를 위해 shape 참조 저장
                        }
                        found_section_on_this_slide = True

                        # Clinical 섹션인 경우 기본 시작점 업데이트
                        if section_type == "clinical":
                            self.body_top = self.section_locations[section_type]["top"]

                # Footer 위치 파악 (전역 설정도 업데이트)
                if "검사기관" in text or "세브란스병원" in text:
                    if shape.top < self.page_height:
                        limit = shape.top - Cm(0.5)
                        if limit < current_slide_limit:
                            current_slide_limit = limit
                        # 하위 호환성을 위해 전역 변수도 업데이트
                        self.body_bottom = current_slide_limit

                # 소제목 삭제 대상 수집
                if found_section_on_this_slide:
                    text_lower = text.lower()
                    is_remove_target = any(k.lower() in text_lower for k in self.config.REMOVE_KEYWORDS)
                    if is_remove_target:
                        shapes_to_remove.append(shape)
                
                # 고지문(Disclaimer) 감지 -> 저장 후 제거 (나중에 동적으로 다시 그림)
                if text.startswith(self.config.MARKER_DISCLAIMER):
                    self.existing_elements['disclaimer'] = copy.deepcopy(shape.element)
                    shapes_to_remove.append(shape)

        # 분석된 슬라이드의 한계선을 ID 기반으로 저장
        self.slide_bottom_limits[slide.slide_id] = current_slide_limit

        # 2. 테이블 요소 분석 및 프로토타입 추출
        for shape in slide.shapes:
            if shape.has_table:
                tbl = shape.table
                try:
                    if len(tbl.rows) > 0:
                        headers = [cell.text_frame.text.strip().upper() for cell in tbl.rows[0].cells]
                        header_str = " ".join(headers)

                        target_key = self._identify_table_type(header_str)

                        if target_key:
                            if target_key not in self.existing_elements['prototypes']:
                                self.existing_elements['prototypes'][target_key] = copy.deepcopy(shape.element)

                            if found_section_on_this_slide:
                                shapes_to_remove.append(shape)

                except Exception as e:
                    print(f"Table analysis warning: {e}")

        # 추출된 요소 제거
        for shape in shapes_to_remove:
            try:
                sp = shape._element
                sp.getparent().remove(sp)
            except ValueError:
                pass

    def _identify_table_type(self, header_str: str) -> str:
        """Config 규칙에 따라 테이블 타입을 식별하는 헬퍼 메서드"""
        for rule in self.config.TABLE_IDENTIFICATION_RULES:
            # Required 조건 체크 (모든 조건 만족해야 함)
            is_match = True
            for req in rule["required"]:
                if isinstance(req, tuple):  # Tuple은 OR 조건
                    if not any(r in header_str for r in req):
                        is_match = False
                        break
                else:  # String은 AND 조건
                    if req not in header_str:
                        is_match = False
                        break

            # Excluded 조건 체크 (하나라도 있으면 탈락)
            if is_match:
                for exc in rule["excluded"]:
                    if exc in header_str:
                        is_match = False
                        break

            if is_match:
                return rule["key"]
        return None


class LayoutContext:
    def __init__(self, prs, analyzer: LayoutAnalyzer, generator=None):
        self.prs = prs
        self.analyzer = analyzer
        self.generator = generator
        self.current_main_title = None
        self.current_slide = prs.slides[0]

        start_loc = analyzer.section_locations.get("clinical")
        if start_loc:
            self.current_slide_index = start_loc["slide_index"]
            self.top = start_loc["top"]
        else:
            self.current_slide_index = 0
            self.top = analyzer.body_top

        self.current_slide = prs.slides[self.current_slide_index]
        self.margin = analyzer.config.MARGIN_LEFT_L3 # Default for content/tables
        self.width = analyzer.config.DEFAULT_WIDTH
        self.config = analyzer.config

    @property
    def bottom_limit(self):
        """현재 보고 있는 슬라이드의 하단 한계선을 동적으로 가져옴"""
        # 슬라이드 ID로 조회
        return self.analyzer.slide_bottom_limits.get(
            self.current_slide.slide_id,
            self.analyzer.config.BODY_BOTTOM_LIMIT
        )

    def move_to_section(self, section_type):
        """특정 섹션의 지정된 위치(템플릿 페이지)로 이동"""
        loc = self.analyzer.section_locations.get(section_type)
        if loc:
            target_slide = self._find_slide_by_id(loc['slide_id'])
            if target_slide:
                self.current_slide = target_slide
                self.top = loc["top"]
                self.current_slide_index = loc["slide_index"]  # 인덱스 업데이트
                return True
        return False

    def check_space(self, height):
        if self.top + height > self.bottom_limit:
            self.add_new_slide()

    def add_new_slide(self):
        # 1. 새 슬라이드 생성 (기본적으로 맨 뒤에 추가됨)
        layout_idx = 6 if len(self.prs.slide_layouts) > 6 else -1
        blank_layout = self.prs.slide_layouts[layout_idx]
        new_slide = self.prs.slides.add_slide(blank_layout)

        # [테두리 복사] 템플릿의 2번째 슬라이드(인덱스 1)에서 배경/테두리 스타일 복사
        if len(self.prs.slides) > 1:
            reference_slide = self.prs.slides[1]
            self._copy_template_style(new_slide, reference_slide)
        
        # [공간 확장] 새로 생성된 슬라이드는 테두리 끝까지 쓸 수 있도록 하단 한계선을 24.5cm로 확장
        self.analyzer.slide_bottom_limits[new_slide.slide_id] = Cm(24.5)

        # 2. 목표 인덱스 계산 (현재 + 1)
        target_index = self.current_slide_index + 1

        # 3. 목표가 맨 뒤가 아니라면 슬라이드 이동
        if target_index < len(self.prs.slides) - 1:
            xml_slides = self.prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[-1])  # 마지막(방금 추가한 것) 제거
            xml_slides.insert(target_index, slides[-1])  # 목표 위치에 삽입
            
            # 새 위치에서 슬라이드 객체 다시 가져오기
            self.current_slide = self.prs.slides[target_index]
        else:
            self.current_slide = new_slide

        # 4. 트래킹 정보 업데이트
        self.current_slide_index = target_index
        self.top = Cm(0.5)

        # [Repeat Main Header] 새 슬라이드에서도 메인 섹션 제목 그리기
        if self.current_main_title and self.generator:
            self.generator._draw_main_section_title(self, self.current_main_title)

    def _copy_template_style(self, target_slide, source_slide):
        """소스 슬라이드에서 배경이나 테두리 같은 정적 도형을 복사"""
        for shape in source_slide.shapes:
            # 텍스트나 테이블이 아닌, 오토쉐이프(도형)나 그림만 복사 대상
            # (테두리는 보통 직사각형 도형으로 구현됨)
            if shape.shape_type in [1, 6]: # 1: AUTO_SHAPE, 6: GROUP
                # 테이블이나 텍스트 프레임이 있는 경우 제외 (단, 단순 도형도 텍스트 프레임 가질 수 있음)
                # 여기서는 '테두리'로 추정되는 크기나 특성을 필터링하면 좋겠지만, 
                # 우선 텍스트가 없는 도형을 복사하는 것으로 단순화
                if not shape.has_table:
                     if not shape.has_text_frame or not shape.text_frame.text.strip():
                         new_el = copy.deepcopy(shape.element)
                         target_slide.shapes._spTree.append(new_el)

    def add_space(self, height):
        self.top += height

    def _find_slide_by_id(self, slide_id):
        for slide in self.prs.slides:
            if slide.slide_id == slide_id:
                return slide
        return None


class NGS_PPT_Generator:
    def __init__(self):
        self.base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        self.template_dir = os.path.join(self.base_dir, "resources")
        self.config = PPTReportConfig()  # 설정 인스턴스

    def _set_cell_border(self, cell, border_color="000000", border_width='12700'):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        for lines in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
            ln = OxmlElement(lines)
            ln.set('w', border_width)
            ln.set('cap', 'flat')
            ln.set('cmpd', 'sng')
            ln.set('algn', 'ctr')
            solidFill = OxmlElement('a:solidFill')
            srgbClr = OxmlElement('a:srgbClr')
            srgbClr.set('val', border_color)
            solidFill.append(srgbClr)
            ln.append(solidFill)
            tcPr.append(ln)
        return cell

    def generate(self, report_data: dict) -> BytesIO:
        panel_type = report_data.get('panel_type', 'GE')
        template_name = "blank_SA_report.pptx" if panel_type == 'SA' else "blank_GE_report.pptx"
        template_path = os.path.join(self.template_dir, template_name)

        if not os.path.exists(template_path):
            raise FileNotFoundError(f"템플릿 파일을 찾을 수 없습니다: {template_path}")

        prs = Presentation(template_path)

        if len(prs.slides) > 0:
            self._fill_clinical_info(prs.slides[0], report_data)
            
            # [Added] 기존 템플릿의 Other Biomarkers 섹션 업데이트
            biomarkers = report_data.get('biomarkers', {})
            self._update_existing_biomarkers(prs.slides[0], biomarkers)

            # [Added] 검사기기 등 Diagnostic Info 업데이트
            self._fill_diagnostic_info(prs, report_data)

        # QC 및 DNA/RNA 정보는 4번째 페이지(Index 3)에 위치함
        if len(prs.slides) > 3:
            self._fill_qc_table(prs.slides[3], report_data)
            self._fill_dna_rna_info(prs.slides[3], report_data)


        analyzer = LayoutAnalyzer(prs)
        self._process_all_variants(prs, report_data, analyzer)
        
        # 코멘트 섹션의 빈 슬라이드(Ghost Page)만 안전하게 제거
        self._remove_ghost_comment_slides(prs)
        
        # 마지막 페이지에 Footer 정보(Tested by, Signed by 등) 기입
        if len(prs.slides) > 0:
            self._fill_footer_info(prs.slides[-1], report_data)

        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output

    def _remove_ghost_comment_slides(self, prs):
        """
        'Comments' 헤더는 있지만 실제 코멘트 내용(본문 박스)이 비어있는 슬라이드를 찾아 삭제합니다.
        (정적 슬라이드인 Method, Gene Content 등은 건드리지 않음)
        """
        slides_to_remove = []
        
        # 헤더 키워드 (부분 일치 허용)
        # '3. Comments', '▣ Comment', '▣ Comment (continued)' 등 모두 커버
        HEADER_KEYWORDS = ["Comment", "Comments"]
        
        for i, slide in enumerate(prs.slides):
            # print(f"DEBUG: Checking Slide {i} for ghost content...")
            has_comment_header = False
            has_meaningful_content = False
            
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                
                # print(f"  - Found Text: '{text[:30]}...' (Len: {len(text)})")

                # 헤더 확인 (유연한 매칭)
                if any(k in text for k in HEADER_KEYWORDS) and len(text) < 40:
                    if "▣" in text or text.startswith("3."): # 3. Comments
                         has_comment_header = True
                         continue
                
                # Footer 확인 (무시)
                # "• 본 검사의 raw data" 혹은 "* 본 검사의 raw data" 등 다양한 불렛 대응
                if "raw data" in text and "보관" in text:
                    continue
                    
                # 그 외 텍스트가 있다면 의미있는 콘텐츠로 간주 (코멘트 본문, 고지문 등)
                has_meaningful_content = True
            
            # 헤더 유무와 상관없이, '의미 있는 본문'이 없으면 Ghost Page로 간주하고 삭제
            # (Header만 있거나, Footer만 있거나, 아예 빈 슬라이드인 경우 모두 포함)
            if not has_meaningful_content:
                slides_to_remove.append(i)
        
        # 역순 삭제
        for idx in sorted(slides_to_remove, reverse=True):
            try:
                # Slide 객체가 아닌 _sldIdLst(XML 요소)에서 rId를 가져와야 함
                slide_id_elem = prs.slides._sldIdLst[idx]
                rId = slide_id_elem.rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[idx]
            except Exception:
                 pass

    def _fill_clinical_info(self, slide, report_data):
        # 1. clinical_info 딕셔너리 추출 (없으면 빈 딕셔너리)
        inner_info = report_data.get('clinical_info', {})
        tables = [shape.table for shape in slide.shapes if shape.has_table]

        for ppt_label, data_key in self.config.CLINICAL_INFO_MAPPING.items():
            # 1순위: clinical_info 내부 검색
            value = inner_info.get(data_key)

            # 2순위: 데이터 루트에서 검색 (대소문자 호환 추가)
            if value is None:
                # 'tmb'로 찾고 없으면 'TMB'로도 찾아봄
                value = report_data.get(data_key) or report_data.get(data_key.upper())

            # 3순위: biomarkers 키 내부 검색 (대소문자 호환 추가)
            if value is None and 'biomarkers' in report_data:
                bio_data = report_data['biomarkers']
                # TMB/MSI가 딕셔너리 형태일 경우 value 추출
                val_obj = bio_data.get(data_key) or bio_data.get(data_key.upper())
                if isinstance(val_obj, dict):
                    value = val_obj.get('value', '')
                else:
                    value = val_obj

            final_value = str(value) if value is not None else ""

            self._search_and_fill_cell_below(tables, ppt_label, final_value, font_size=Pt(8))

    def _update_existing_biomarkers(self, slide, biomarkers):
        """
        기존 템플릿(Slide 1)에 존재하는 'Other Biomarkers' 섹션을 찾아 업데이트합니다.
        
        Logic:
        1. 슬라이드 내의 모든 Shape를 순회하며:
           - "- Other Biomarkers" 텍스트를 가진 Shape를 찾아 제목 업데이트 (Status Suffix 추가)
           - TMB/MSI 헤더를 가진 테이블 Shape를 찾아 데이터 셀 업데이트 (값 + 단위, Partial Styling)
        """
        if not biomarkers:
            return

        tmb_data = biomarkers.get('TMB', {})
        msi_data = biomarkers.get('MSI', {})
        
        # 딕셔너리가 아닌 경우(구버전 호환) 처리
        if not isinstance(tmb_data, dict): tmb_data = {'value': str(tmb_data), 'status': '', 'unit': ''}
        if not isinstance(msi_data, dict): msi_data = {'value': str(msi_data), 'status': '', 'unit': ''}

        tmb_val = tmb_data.get('value', '')
        tmb_unit = tmb_data.get('unit', '')
        
        msi_val = msi_data.get('value', '')
        msi_unit = msi_data.get('unit', '')
        
        # Status Check
        is_tmb_high = 'high' in str(tmb_data.get('status', '')).lower()
        is_msi_high = 'high' in str(msi_data.get('status', '')).lower()

        # Suffix Calculation
        suffix = ""
        if is_tmb_high and is_msi_high:
            suffix = ": TMB-High and MSI-High"
        elif is_tmb_high:
            suffix = ": TMB-High"
        elif is_msi_high:
            suffix = ": MSI-High"

        # Shape Iteration
        for shape in slide.shapes:
            # 1. Update Title
            if shape.has_text_frame:
                if "- Other Biomarkers" in shape.text_frame.text:
                    p = shape.text_frame.paragraphs[0]
                    
                    # [Fix] 텍스트 박스 너비를 넓게 조정하여 줄바꿈 방지
                    # 기존 너비(약 4.6cm)로는 "TMB-High and MSI-High" 추가 시 줄바꿈됨.
                    shape.width = Cm(17.0)
                    shape.text_frame.word_wrap = False
                    
                    # Suffix (Red/Bold) 추가
                    if suffix:
                        run = p.add_run()
                        run.text = suffix
                        run.font.name = self.config.FONT_NAME
                        run.font.size = self.config.FONT_SIZE_HEADER
                        run.font.bold = True
                        run.font.color.rgb = self.config.COLOR_RED

            # 2. Update Table
            if shape.has_table:
                tbl = shape.table
                # 헤더 확인 (TMB/MSI가 포함되어 있는지)
                # 보통 첫 번째 행(Header)의 내용을 검사
                try:
                    header_row_text = "".join([cell.text_frame.text for cell in tbl.rows[0].cells])
                    if "Tumor Mutation Burden" in header_row_text and "Microsatellite Instability" in header_row_text:
                        # 타겟 테이블 발견
                        
                        # TMB Cell (Row 1, Col 0)
                        # 기존에 "/Megabase" 같은 단위 텍스트가 있을 수 있으므로 덮어쓰기
                        self._fill_biomarker_cell(tbl.cell(1, 0), tmb_val, tmb_unit, is_tmb_high)

                        # MSI Cell (Row 1, Col 1)
                        self._fill_biomarker_cell(tbl.cell(1, 1), msi_val, msi_unit, is_msi_high)
                except Exception:
                    continue
        
    def _set_cell_style_simple(self, cell, is_header=False):
        """테이블 셀 스타일 유틸리티 (헤더용)"""
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        
        if p.runs:
            font = p.runs[0].font
            font.name = self.config.FONT_NAME
            font.size = Pt(10)
            font.bold = True
            font.color.rgb = self.config.COLOR_BLACK

        if is_header:
            from pptx.dml.color import RGBColor
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(220, 230, 241) # 연한 파랑/회색 계열

        self._set_cell_border(cell)

    def _fill_biomarker_cell(self, cell, value, unit, is_high):
        """바이오마커 셀 채우기 (값 부분만 강조)"""
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = "" # 기존 텍스트 클리어

        # Value
        run_val = p.add_run()
        run_val.text = str(value)
        run_val.font.name = self.config.FONT_NAME
        run_val.font.size = Pt(8) # [Changed] 10pt -> 8pt
        
        if is_high:
            run_val.font.bold = True
            run_val.font.color.rgb = self.config.COLOR_RED
        else:
            run_val.font.bold = False
            run_val.font.color.rgb = self.config.COLOR_BLACK
            
        # Unit (공백 추가)
        if unit:
            run_unit = p.add_run()
            run_unit.text = f" {str(unit)}"
            run_unit.font.name = self.config.FONT_NAME
            run_unit.font.size = Pt(8) # [Changed] 10pt -> 8pt
            run_unit.font.bold = False
            run_unit.font.color.rgb = self.config.COLOR_BLACK

        self._set_cell_border(cell)

    def _fill_qc_table(self, slide, report_data):
        qc_data = report_data.get('qc', {})
        if not qc_data: return

        rows_data = qc_data.get('data', [])
        if not rows_data: return

        qc_map = {}
        for row in rows_data:
            if len(row) >= 2:
                key = str(row[0]).strip().replace(" ", "").lower()
                val = str(row[-1])
                qc_map[key] = val
        
        tables = [shape.table for shape in slide.shapes if shape.has_table]
        for table in tables:
            for row in table.rows:
                if len(row.cells) > 0:
                    first_text = row.cells[0].text_frame.text.strip().replace(" ", "").lower()
                    
                    matched_val = None
                    for k, v in qc_map.items():
                        # 부분 일치 허용 (예: 데이터 키 'pct_pf_reads'가 테이블 텍스트 'pct_pf_reads(%)'에 포함되는지 확인)
                        # 혹은 반대 경우도 고려
                        if k in first_text or first_text in k:
                            matched_val = v
                            break
                    
                    if matched_val:
                        # 마지막 컬럼에 채우기
                        last_idx = len(row.cells) - 1
                        self._set_cell_text_preserving_style(row.cells[last_idx], matched_val, is_bold=True)

    def _fill_dna_rna_info(self, slide, report_data):
        drna = report_data.get('drna_qubit', {})
        if not drna: return
        
        dna_val = drna.get('DNA')
        rna_val = drna.get('RNA')
        
        t_dna = "DNA (ng/ul)"
        t_rna = "RNA (ng/ul)"

        for shape in slide.shapes:
            if shape.has_text_frame:
                self._process_text_frame_for_drna(shape.text_frame, t_dna, t_rna, dna_val, rna_val)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        self._process_text_frame_for_drna(cell.text_frame, t_dna, t_rna, dna_val, rna_val)

    def _process_text_frame_for_drna(self, text_frame, t_dna, t_rna, v_dna, v_rna):
        if t_dna not in text_frame.text and t_rna not in text_frame.text:
            return
            
        for p in text_frame.paragraphs:
            txt = p.text
            modified = False
            
            # DNA 채우기
            if v_dna and t_dna in txt:
                idx = txt.find(t_dna)
                c_idx = txt.find(":", idx)
                if c_idx != -1:
                    # 중복 삽입 방지 로직이 없으므로, 템플릿이 깨끗하다고 가정하고 삽입
                    prefix = txt[:c_idx+1]
                    suffix = txt[c_idx+1:]
                    # 값 삽입 (한 칸 띄움)
                    txt = f"{prefix} {v_dna}{suffix}"
                    modified = True
            
            # RNA 채우기 (DNA 처리된 텍스트에서 검색)
            if v_rna and t_rna in txt:
                idx = txt.find(t_rna) # 뒤쪽에 위치한 RNA 찾기
                c_idx = txt.find(":", idx)
                if c_idx != -1:
                    prefix = txt[:c_idx+1]
                    suffix = txt[c_idx+1:]
                    txt = f"{prefix} {v_rna}{suffix}"
                    modified = True
            
            if modified:
                p.text = txt
                # 스타일(폰트) 유지 노력
                if p.runs:
                    p.runs[0].font.name = self.config.FONT_NAME

    def _extract_highlight_keywords(self, report_data):
        """Report Data에서 코멘트 강조를 위한 키워드(유전자명 등)를 추출합니다."""
        keywords = set()
        target_headers = ['GENE', 'MUTATION', 'AA CHANGE', 'FUSION', 'SPLICE', 'COPY', 'BREAKPOINT']
        
        # 제외할 일반적인 단어 리스트 (단독으로 Bold 처리되지 않도록)
        BLACKLIST = {"mutation", "fusion", "amplification", "deletion", "duplication", 
                     "insertion", "splice", "variant", "copy", "loss", "gain", 
                     "rearrangement", "skipping", "indel", "snv", "cnv", "none"}
        
        # 1. Variant Sections 순회
        for section_config in self.config.VARIANT_SECTIONS:
            key = section_config['key']
            section_data = report_data.get(key, {})
            rows = section_data.get('data', [])
            headers = section_data.get('headers', [])
            
            if not rows or not headers:
                continue
            
            # 헤더 인덱스 매핑 (대소문자 무관하게 처리)
            indices = {}
            for h_idx, h_val in enumerate(headers):
                h_upper = str(h_val).upper().strip()
                if h_upper in target_headers:
                    indices[h_upper] = h_idx
            
            for row in rows:
                gene = None
                
                # GENE
                if 'GENE' in indices and len(row) > indices['GENE']:
                    val = str(row[indices['GENE']]).strip()
                    if val and val.lower() != 'none':
                        gene = val
                        keywords.add(gene)
                        # 편의상 예상되는 조합형 키워드 추가
                        keywords.add(f"{gene} mutation") 
                        keywords.add(f"{gene} amplification")
                        keywords.add(f"{gene} fusion")
                
                # 그 외 컬럼들 (MUTATION, AA CHANGE, FUSION 등)
                for header_key in ['MUTATION', 'AA CHANGE', 'FUSION', 'SPLICE']:
                    if header_key in indices and len(row) > indices[header_key]:
                        val = str(row[indices[header_key]]).strip()
                        if val and val.lower() not in BLACKLIST:
                            # 1) 단독 키워드 추가 (블랙리스트에 없으면)
                            keywords.add(val)
                        
                        # 2) Gene 결합형 키워드 추가 (항상 추가)
                        if gene and val and val.lower() != 'none':
                            # 중복 공백 제거
                            composite = f"{gene} {val}"
                            keywords.add(composite)

        # 2. Failed Gene
        failed_gene = report_data.get('failed_gene')
        if failed_gene and str(failed_gene).lower() != 'none':
             keywords.add(str(failed_gene))
        
        # 3. 추가 기본 키워드 (사용자 요청 사례 등)
        keywords.add("MET exon14 skipping")
        
        # 최종 필터링: 너무 짧은 단어(2글자 이하) 등 제거
        final_keywords = [k for k in keywords if len(k) > 2]
        return list(set(final_keywords))

    def _process_all_variants(self, prs, report_data, analyzer):
        layout = LayoutContext(prs, analyzer, generator=self)
        current_section_group = None

        for section_config in self.config.VARIANT_SECTIONS:
            key = section_config['key']
            title = section_config['title']
            style_type = section_config['type']
            prototype_key = section_config['prototype_key']

            # 스타일 설정 가져오기 (clinical: Red/Bold, unknown: Black/Normal)
            style_props = self.config.STYLES.get(style_type, self.config.STYLES["unknown"])

            if current_section_group != style_type:
                current_section_group = style_type
                
                # [리팩토링된 흐름 로직]
                # 템플릿의 앵커로 이동할지, 아니면 현재 흐름을 유지할지 결정
                target_loc = analyzer.section_locations.get(style_type)
                should_move = False
                
                if target_loc:
                    target_slide_idx = target_loc["slide_index"]
                    target_slide_id = target_loc["slide_id"]

                    # 1. 이미 현재 슬라이드가 목표 슬라이드인 경우 (예: Clinical 섹션 시작)
                    # 별도 이동 없이 앵커를 사용하도록 허용
                    if layout.current_slide.slide_id == target_slide_id:
                        should_move = True
                    # 2. 목표 슬라이드를 아직 지나지 않았을 때만 이동
                    # 엄격한 부등호(<) 사용: 오버플로우로 인해 현재 인덱스가 목표 인덱스와 같아진 경우 점프 방지
                    elif layout.current_slide_index < target_slide_idx:
                        should_move = True
                
                found_anchor = False
                if should_move:
                    found_anchor = layout.move_to_section(style_type)
                else:
                    # 이동하지 않기로 결정했다면(오버플로우 등), 
                    # 중복 방지를 위해 원래 템플릿 슬라이드의 제목을 제거해야 함
                    if target_loc and "title_shape" in target_loc:
                        try:
                            sp = target_loc["title_shape"]
                            if sp.element.getparent() is not None:
                                sp.element.getparent().remove(sp.element)
                        except Exception as e:
                            print(f"Warning: Failed to remove unused title shape: {e}")
                
                # 이동하지 않았거나(오버플로우), 앵커를 못 찾았다면 직접 타이틀 그리기
                if not found_anchor:
                    main_title = self.config.SECTION_START_MARKERS.get(style_type)
                    if main_title:
                        # 기본값이 12pt Black으로 변경되었으므로 별도 파라미터 불필요
                        self._draw_main_section_title(layout, main_title)
                
                # [Tracking] 현재 메인 타이틀 업데이트 (페이지 분할 시 반복용)
                layout.current_main_title = self.config.SECTION_START_MARKERS.get(style_type)

            # [Added] LR-BRCA (Clinical) 데이터가 없으면 섹션 자체를 스킵 (PDF와 동작 일치)
            if key == 'lr_brca_clinical' and not section_data.get('data'):
                continue

            section_data = report_data.get(key, {})
            rows = section_data.get('data', [])
            headers = section_data.get('headers', [])
            highlight_val = section_data.get('highlight', [])

            prototype_xml = analyzer.existing_elements["prototypes"].get(prototype_key)

            if rows and len(rows) > 0:
                header_height = Cm(0.8)
                min_table_height = Cm(0.8)
                required_height = header_height + min_table_height

                # [Changed] Orphan Header Prevention
                # 헤더(0.8) + 최소 행 높이(약 1.0) 정도의 여유 공간 확인
                # 부족하면 아예 새 슬라이드로 넘김
                orphan_threshold = header_height + Cm(1.2)
                
                if layout.top + orphan_threshold > layout.bottom_limit:
                    layout.add_new_slide()

                # 먼저 테이블 분할 여부 계산
                if prototype_xml is not None:
                    total_pages = self._calculate_table_pages(layout, len(rows))
                else:
                    total_pages = 1  # scratch 테이블은 분할 미지원
                
                # 분할 시 첫 페이지 제목에 (1/N) 추가
                display_title = f"{title} (1/{total_pages})" if total_pages > 1 else title
                
                # [Changed] VUS 여부 확인 (highlight plain text 처리를 위해)
                # style_type 변수를 사용 (section_data에는 type 정보가 없을 수 있음)
                is_clinical = (style_type == 'clinical')
                
                self._render_section_header(layout, display_title, highlight_data=highlight_val, is_clinical=is_clinical)
                
                if prototype_xml is not None:
                    # 테이블 분할 시 제목과 페이지 번호 표시를 위한 정보 전달
                    if total_pages > 1:
                        pagination_info = {
                            "title": title,
                            "current_page": 1,
                            "total_pages": total_pages,
                            "is_first": True
                        }
                    else:
                        pagination_info = None
                    
                    self._render_table_using_prototype(
                        layout, prototype_xml, rows, style_props,
                        margin_left=self.config.MARGIN_LEFT_L3,
                        pagination_info=pagination_info
                    )
                else:
                    self._render_table_from_scratch(layout, headers, rows, margin_left=self.config.MARGIN_LEFT_L3)
            else:
                self._render_section_header(layout, title, is_none=True)

            layout.add_space(self.config.SPACE_SECTION)

        # [Safety Reset] Variants 섹션 종료 후, 후속 섹션(Failed Gene 등)에 이전 타이틀이 반복되지 않도록 초기화
        layout.current_main_title = None

        # 3. Failed gene 섹션 (Variants 처리 루프 종료 후)
        failed_gene = report_data.get('failed_gene')
        self._draw_failed_gene(layout, failed_gene)

        # 5. Other Biomarkers
        # [Refactoring] 기존 템플릿(Slide 1)에 있는 섹션을 업데이트하므로 여기서는 그리지 않음.
        # biomarkers = report_data.get('biomarkers', {})
        # self._draw_biomarkers(layout, biomarkers)

        # Comments 섹션
        comments = report_data.get('comments', [])
        highlight_keywords = self._extract_highlight_keywords(report_data)
        self._draw_comments(layout, comments, highlight_keywords)

    def _draw_main_section_title(self, layout, text, font_size=None, color=None):
        height = Cm(0.6)
        layout.check_space(height)
        tb = layout.current_slide.shapes.add_textbox(self.config.MARGIN_LEFT_L1, layout.top, layout.width, height)
        p = tb.text_frame.paragraphs[0]
        
        # 기본값 설정
        final_size = font_size if font_size else self.config.FONT_SIZE_TITLE
        
        # [Color Logic Improvement for Header Repetition]
        # "1. ", "2. " 등 번호 부분은 검정색, 뒷부분은 섹션 성격에 따라 색상 적용
        # 명시적 color 파라미터가 있으면 그것을 따름 (기존 호환성)
        
        import re
        match = re.match(r"^(\d+\.)\s*(.*)", text)
        
        if match and not color:
            number_part = match.group(1)
            title_part = match.group(2)
            
            # 색상 결정: "Unknown"이 포함되면 검정, 아니면(Clinical) 빨강
            # (Unknown Variants 제목은 보통 검정색, Clinical은 빨강색)
            is_unknown = "Unknown" in text or "unknown" in text
            title_color = self.config.COLOR_BLACK if is_unknown else self.config.COLOR_RED
            
            # 1. Number Part (Always Black)
            self._set_run_style(p.add_run(), f"{number_part} ", is_bold=True,
                                font_size=final_size, color=self.config.COLOR_BLACK)
            
            # 2. Title Part (Context Color)
            self._set_run_style(p.add_run(), title_part, is_bold=True,
                                font_size=final_size, color=title_color)
        else:
            # 기존 로직 (전체 검정 또는 지정색)
            final_color = color if color else self.config.COLOR_BLACK 
            self._set_run_style(p.add_run(), text, is_bold=True,
                                font_size=final_size,
                                color=final_color)

        layout.add_space(height + self.config.SPACE_TITLE_BOTTOM)

    def _render_section_header(self, layout, title, highlight_data=None, is_none=False, is_clinical=True):
        height = Cm(0.8)
        layout.check_space(height)
        
        # 텍스트 박스 너비를 명시적으로 제한하여 우측 테두리 안쪽에서 안전하게 자동 줄바꿈(word wrap) 유도
        box_width = layout.width - self.config.MARGIN_LEFT_L2 - Cm(0.5)
        tb = layout.current_slide.shapes.add_textbox(self.config.MARGIN_LEFT_L2, layout.top, box_width, height)
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]

        # 1. 제목
        self._set_run_style(p.add_run(), f"- {title}", is_bold=True,
                            font_size=self.config.FONT_SIZE_HEADER,
                            color=self.config.COLOR_BLACK)

        # 2. 상태값
        if is_none:
            self._set_run_style(p.add_run(), ": None",
                                font_size=self.config.FONT_SIZE_HEADER,
                                color=self.config.COLOR_BLACK)
        elif highlight_data:
            # [Changed] VUS(is_clinical=False)인 경우: Plain Text (Black, Normal, No Italic)
            # VCS(is_clinical=True)인 경우: 기존 로직 (Red, Bold, Gene Italic)
            
            # 구분자 색상/스타일 결정
            sep_color = self.config.COLOR_RED if is_clinical else self.config.COLOR_BLACK
            sep_bold = True if is_clinical else False # VUS는 구분자도 Normal
            
            # ": " 찍기
            self._set_run_style(p.add_run(), ": ", is_bold=sep_bold,
                                font_size=self.config.FONT_SIZE_HEADER,
                                color=sep_color)

            # 데이터 순회
            for segment in highlight_data:
                text = segment['text']
                style = segment.get('style', 'normal')

                if is_clinical:
                    # VCS: Red, Bold, Italic(if gene)
                    is_italic = (style == 'italic')
                    self._set_run_style(p.add_run(), text, is_bold=True,
                                        font_size=self.config.FONT_SIZE_HEADER,
                                        color=self.config.COLOR_RED,
                                        italic=is_italic)
                else:
                    # VUS: Black, Normal, No Italic (Plain Text)
                    # style 정보(italic 등) 무시하고 강제 평문 출력
                    self._set_run_style(p.add_run(), text, is_bold=False,
                                        font_size=self.config.FONT_SIZE_HEADER,
                                        color=self.config.COLOR_BLACK,
                                        italic=False) 

        # 동적 높이 계산 (Dynamic Height Calculation with constraints)
        wrapped_lines = 0
        if highlight_data:
            total_chars = sum(len(segment['text']) for segment in highlight_data)
            # 대략적인 한 줄 글자 수 (12pt 맑은 고딕, 17cm 너비 기준 보수적 측정치)
            chars_per_line = 80 
            if total_chars > chars_per_line:
                wrapped_lines = int((total_chars - 1) / chars_per_line)
                
        # 최대 3줄(기본 1줄 + 추가 2줄)까지만 공간 할당을 허용
        capped_wrapped_lines = min(wrapped_lines, 2)
        
        # 줄바꿈 1개당 0.5cm씩 추가
        dynamic_added_height = Cm(0.5) * capped_wrapped_lines
        layout.add_space(height + dynamic_added_height)

    def _set_run_style(self, run, text, is_bold=False, font_size=None, color=None, italic=False):
        """텍스트 스타일 적용 헬퍼 메서드"""
        run.text = text
        run.font.name = self.config.FONT_NAME
        if font_size: run.font.size = font_size
        if color: run.font.color.rgb = color
        run.font.bold = is_bold
        run.font.italic = italic

    def _calculate_table_pages(self, layout, total_rows):
        """테이블이 몇 페이지에 걸쳐 분할될지 미리 계산합니다."""
        row_height = Cm(0.7)
        header_height = Cm(0.8)
        page_count = 0
        remaining_rows = total_rows
        
        # 첫 페이지 가용 공간
        available_height = layout.bottom_limit - layout.top
        if available_height < (header_height + row_height):
            available_height = layout.bottom_limit - self.config.BODY_TOP_START
        
        while remaining_rows > 0:
            page_count += 1
            max_rows = int((available_height - header_height) / row_height)
            remaining_rows -= max(1, max_rows)  # 최소 1행은 들어감
            # 이후 페이지는 새 슬라이드 기준
            available_height = layout.bottom_limit - self.config.BODY_TOP_START
        
        return max(1, page_count)

    def _render_table_using_prototype(self, layout, prototype_xml, rows, style_props, 
                                       margin_left=None, pagination_info=None):
        """프로토타입 XML을 사용하여 테이블을 그립니다.
        
        Args:
            pagination_info: 분할 시 페이지 번호 표시를 위한 정보 dict
                - title: 섹션 제목 (예: "SNVs & Indels")
                - current_page: 현재 페이지 번호
                - total_pages: 전체 페이지 수
                - is_first: 첫 페이지 여부 (첫 페이지는 별도 렌더링되므로 제목 스킵)
        """
        final_margin = margin_left if margin_left else layout.margin
        row_height = Cm(0.7)
        header_height = Cm(0.8)

        available_height = layout.bottom_limit - layout.top
        
        # [DEBUG]
        from pptx.util import Length
        avail_cm = Length(available_height).cm if isinstance(available_height, int) else available_height.cm
        bottom_cm = Length(layout.bottom_limit).cm if isinstance(layout.bottom_limit, int) else layout.bottom_limit.cm
        top_cm = Length(layout.top).cm if isinstance(layout.top, int) else layout.top.cm
        
        print(f"[DEBUG Layout] Table rendering start. Available Height: {avail_cm:.2f}cm (Bottom: {bottom_cm:.2f}cm - Top: {top_cm:.2f}cm)")

        if available_height < (header_height + row_height):
            print(f"[DEBUG Layout] Not enough space for header+1row. Adding new slide.")
            layout.add_new_slide()
            available_height = layout.bottom_limit - layout.top
            avail_cm = Length(available_height).cm if isinstance(available_height, int) else available_height.cm

        max_rows = int((available_height - header_height) / row_height)
        print(f"[DEBUG Layout] Calculating max_rows: ({avail_cm:.2f} - {header_height.cm:.2f}) / {row_height.cm:.2f} = {max_rows}")

        # 페이지 정보 초기화 (첫 호출 시)
        if pagination_info is None:
            total_pages = self._calculate_table_pages(layout, len(rows))
            if total_pages > 1:
                pagination_info = {
                    "title": None,  # 호출자가 설정
                    "current_page": 1,
                    "total_pages": total_pages,
                    "is_first": True
                }

        if max_rows >= len(rows):
            # 마지막 페이지 또는 분할 없음
            self._insert_cloned_table(layout, prototype_xml, rows, style_props, final_margin)
        else:
            current_batch = rows[:max_rows]
            next_batch = rows[max_rows:]
            self._insert_cloned_table(layout, prototype_xml, current_batch, style_props, final_margin)
            layout.add_new_slide()
            
            # 다음 페이지 정보 업데이트 및 제목 렌더링
            if pagination_info:
                next_page = pagination_info["current_page"] + 1
                total = pagination_info["total_pages"]
                title = pagination_info.get("title")
                
                # 제목이 있으면 (n/N) 형식으로 렌더링
                if title:
                    title_with_page = f"{title} ({next_page}/{total})"
                    self._render_section_header(layout, title_with_page)
                
                next_info = {
                    "title": title,
                    "current_page": next_page,
                    "total_pages": total,
                    "is_first": False
                }
            else:
                next_info = None
                
            self._render_table_using_prototype(layout, prototype_xml, next_batch, style_props, 
                                                margin_left=final_margin, pagination_info=next_info)

    def _insert_cloned_table(self, layout, prototype_xml, rows, style_props, margin_left):
        new_tbl_element = copy.deepcopy(prototype_xml)
        layout.current_slide.shapes._spTree.insert_element_before(new_tbl_element, 'p:extLst')

        table_shape = layout.current_slide.shapes[-1]
        table_shape.top = int(layout.top)
        table_shape.left = int(margin_left)
        table = table_shape.table

        current_rows = len(table.rows)
        needed_rows = len(rows) + 1

        if needed_rows > current_rows:
            for _ in range(needed_rows - current_rows):
                self._duplicate_last_row(table)

        for r_idx, row_data in enumerate(rows):
            target_row = table.rows[r_idx + 1]
            for c_idx, val in enumerate(row_data):
                if c_idx < len(target_row.cells):
                    self._set_cell_text_preserving_style(
                        target_row.cells[c_idx],
                        str(val),
                        is_bold=style_props['bold'],
                        font_color=style_props['color'],
                        font_size=Pt(8) # [Changed] Explicitly set to 8pt for variants
                    )

        table_height = sum([row.height for row in table.rows])
        layout.add_space(table_height + self.config.SPACE_TABLE_BOTTOM)

    def _set_cell_text_preserving_style(self, cell, text, is_bold=False, font_color=None, font_size=None):
        """텍스트 입력 전 기존 내용을 초기화하여 중복/깨짐 방지"""
        if not cell.text_frame.paragraphs:
            p = cell.text_frame.add_paragraph()
        else:
            p = cell.text_frame.paragraphs[0]

        p.text = ""

        p.alignment = PP_ALIGN.CENTER

        run = p.add_run()
        run.font.name = self.config.FONT_NAME
        
        # [Changed] font_size 파라미터가 있으면 사용, 없으면 Default Body Size (9pt)
        if font_size:
            run.font.size = font_size
        else:
            run.font.size = self.config.FONT_SIZE_BODY

        run.font.bold = is_bold
        if font_color: run.font.color.rgb = font_color
        run.text = text

    def _duplicate_last_row(self, table):
        new_row = copy.deepcopy(table._tbl.tr_lst[-1])
        for tc in new_row.tc_lst:
            txBody = tc.find('{http://schemas.openxmlformats.org/drawingml/2006/main}txBody')
            if txBody is not None:
                p_list = txBody.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}p')
                for p in p_list:
                    for child in list(p):
                        p.remove(child)
        table._tbl.append(new_row)

    def _render_table_from_scratch(self, layout, headers, rows, margin_left=None):
        final_margin = margin_left if margin_left else layout.margin
        rows_count = len(rows) + 1
        cols_count = len(headers)
        table_height = Cm(0.8 * rows_count)

        shape = layout.current_slide.shapes.add_table(
            rows_count, cols_count, final_margin, layout.top, layout.width, table_height
        )
        table = shape.table

        for idx, h in enumerate(headers):
            cell = table.cell(0, idx)
            cell.text = str(h)
            self._set_cell_border(cell)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.config.COLOR_GRAY_BG

        for r, row_data in enumerate(rows):
            for c, val in enumerate(row_data):
                cell = table.cell(r + 1, c)
                cell.text = str(val)
                self._set_cell_border(cell)

        layout.add_space(table_height + self.config.SPACE_TABLE_BOTTOM)

    def _search_and_fill_cell_below(self, tables, target_label, value, font_size=None):
        clean_target = target_label.replace(" ", "").lower()

        for table in tables:
            for r_idx, row in enumerate(table.rows):
                for c_idx, cell in enumerate(row.cells):
                    cell_text = cell.text_frame.text.replace(" ", "").lower()

                    if clean_target in cell_text:
                        try:
                            if r_idx + 1 < len(table.rows):
                                target_cell = table.cell(r_idx + 1, c_idx)

                                # 기존 텍스트(단위) 가져오기 (예: "/Megabase")
                                original_text = target_cell.text_frame.text.strip()
                                
                                # 입력할 값 문자열 정리
                                val_str = str(value).strip() if value is not None else ""

                                if val_str:
                                    # 1. 데이터(val_str) 안에 이미 단위(original_text)가 들어있는 경우
                                    if original_text and original_text in val_str:
                                        final_text = val_str

                                    # 2. 데이터에 단위가 없는 경우 (숫자만 있는 경우)
                                    elif original_text:
                                        final_text = f"{val_str} {original_text}"

                                    # 3. 기존 단위가 아예 없는 셀인 경우
                                    else:
                                        final_text = val_str
                                else:
                                    # 값이 없으면 기존 단위 유지
                                    final_text = original_text

                                # 스타일 적용하여 입력
                                self._set_cell_text_preserving_style(
                                    target_cell,
                                    final_text,
                                    is_bold=True,
                                    font_size=font_size
                                )
                                return
                        except Exception as e:
                            print(f"Warning: Failed to fill {target_label} - {e}")

    def _draw_failed_gene(self, layout, failed_gene_text):
        """3. Failed gene 섹션을 그립니다."""
        # 제목이자 내용이 한 줄에 들어가는 형태: "3. Failed gene: {내용}"
        height = Cm(1.0)
        layout.check_space(height)
        
        tb = layout.current_slide.shapes.add_textbox(self.config.MARGIN_LEFT_L1, layout.top, layout.width, height)
        p = tb.text_frame.paragraphs[0]
        
        # "3. Failed gene" 부분 (Unknown과 동일하게 12pt Black 적용)ㄴ
        self._set_run_style(p.add_run(), "3. Failed gene", is_bold=True,
                            font_size=self.config.FONT_SIZE_TITLE, color=self.config.COLOR_BLACK)
        
        # ": " 및 내용 부분
        content = f": {failed_gene_text}" if failed_gene_text else ": None"
        self._set_run_style(p.add_run(), content, is_bold=False,
                            font_size=self.config.FONT_SIZE_HEADER, color=self.config.COLOR_BLACK)
        
        layout.add_space(height + self.config.SPACE_SECTION)


    def _draw_comments(self, layout, comments, highlight_keywords=None):
        """Comments 섹션과 하단 고지문을 통합하여 그립니다."""
        # 상수 정의 (총 페이지 계산에 필요)
        LINE_HEIGHT = Cm(0.6) # [Changed] 1.5 Spacing 고려하여 높이 증가 (0.4 -> 0.6)
        CHARS_PER_LINE = 90
        FOOTER_TOP = Cm(23.49)
        BODY_BOTTOM_LIMIT = FOOTER_TOP - Cm(0.5)
        if layout.bottom_limit < BODY_BOTTOM_LIMIT:
            BODY_BOTTOM_LIMIT = layout.bottom_limit
        DISCLAIMER_MAIN_HEIGHT = Cm(3.0)
        
        # 총 페이지 수 미리 계산
        if not comments:
            comments_list = []
        elif isinstance(comments, str):
            comments_list = [comments]
        else:
            comments_list = comments
            
        total_pages = self._calculate_comment_pages(
            layout, comments_list, BODY_BOTTOM_LIMIT, LINE_HEIGHT, CHARS_PER_LINE, DISCLAIMER_MAIN_HEIGHT
        )
        
        # 1. 첫 페이지 헤더 그리기 (분할 시 (1/N) 형식)
        page_info = (1, total_pages) if total_pages > 1 else None
        self._draw_comment_header(layout, page_info=page_info)

        # 2. 통합 콘텐츠 그리기 (코멘트 + 고지문)
        self._draw_merged_content(layout, comments, highlight_keywords)

    def _draw_comment_header(self, layout, page_info=None):
        """Comments 헤더를 그립니다.
        
        Args:
            page_info: (current_page, total_pages) 튜플. None이면 번호 없이 그림.
        """
        header_height = Cm(1.0)
        layout.check_space(header_height)
        
        # 박스와 동일한 좌측 여백 계산 (페이지 중앙 정렬)
        BOX_WIDTH = Cm(17.55)
        slide_width = layout.prs.slide_width
        left_position = (slide_width - BOX_WIDTH) / 2
        
        tb_header = layout.current_slide.shapes.add_textbox(left_position, layout.top, layout.width, header_height)
        p_header = tb_header.text_frame.paragraphs[0]
        
        # 페이지 번호 형식: 분할 시 (1/N), (2/N) 등
        if page_info and page_info[1] > 1:
            title_text = f"▣ Comment ({page_info[0]}/{page_info[1]})"
        else:
            title_text = "▣ Comment"
        
        self._set_run_style(p_header.add_run(), title_text, is_bold=True, 
                            font_size=self.config.FONT_SIZE_TITLE, color=self.config.COLOR_BLACK)
        layout.add_space(header_height)

    def _draw_merged_content(self, layout, comments, highlight_keywords):
        if not comments:
            comments = [] 

        if isinstance(comments, str):
            comments = [comments]

        # 상수 정의
        LINE_HEIGHT = Cm(0.6) # [Changed] 1.5 Spacing 고려하여 높이 증가
        CHARS_PER_LINE = 90
        BOX_WIDTH = Cm(17.55)
        
        # 하단 Footer (Raw Data 등) 위치
        FOOTER_TOP = Cm(23.49)
        BODY_BOTTOM_LIMIT = FOOTER_TOP - Cm(0.5)
        
        if layout.bottom_limit < BODY_BOTTOM_LIMIT:
             BODY_BOTTOM_LIMIT = layout.bottom_limit
        
        # Main Disclaimer
        disclaimer_main = (
            "*본 기관의 유전자 정보 검색 및 해석은 dbSNP, COSMIC, ClinVar, c-bioportal 등의 유전자 정보검색 사이트를 참고로 하고 있습니다. "
            "또한, 발견된 유전자 변이의 임상적 의미에 대하여, 아래 참고 문헌에 기반한 4단계 시스템 (Tier I: Strong clinical significance, "
            "Tier II: Potential Clinical significance, Tier III: Unknown Clinical significance, Tier IV: Benign or Likely Benign) 을 "
            "적용하여 보고하고 있습니다. 또한 normal cell sequencing이 이루어지지 않아 일부 변이의 경우 germline polymorphism의 가능성을 배제할 수 없습니다.\n"
            "Standards and Guidelines for the Interpretation and Reporting of Sequence Variants in Cancer. A Joint Consensus Recommendation "
            "of the Association for Molecular Pathology, American Society of Clinical Oncology and College of American Pathologists. "
            "J Mol Diagn. 2017; 19(1):4-23."
        )
        DISCLAIMER_MAIN_HEIGHT = Cm(3.0)
        
        # 총 페이지 수 미리 계산
        total_pages = self._calculate_comment_pages(
            layout, comments, BODY_BOTTOM_LIMIT, LINE_HEIGHT, CHARS_PER_LINE, DISCLAIMER_MAIN_HEIGHT
        )
        current_page = 1
        
        current_batch = []
        current_batch_height = 0 
        
        # 코멘트 배치 루프
        for idx, comment in enumerate(comments):
            text_len = len(comment)
            lines = int((text_len / CHARS_PER_LINE) + 1)
            est_height = (lines * LINE_HEIGHT) + Cm(0.2)
            
            # 공간 체크 (Footer 영역 침범 확인)
            if (layout.top + current_batch_height + est_height) > BODY_BOTTOM_LIMIT:
                 # 넘치면 현재 배치 그리기
                 self._render_box(layout, current_batch, current_batch_height, BOX_WIDTH, highlight_keywords, main_disclaimer=None)
                 self._draw_footer_info(layout, FOOTER_TOP)
                 
                 # 다음 페이지 이동
                 layout.add_new_slide()
                 current_page += 1
                 self._draw_comment_header(layout, page_info=(current_page, total_pages))
                 
                 current_batch = []
                 current_batch_height = 0
            
            current_batch.append(comment)
            current_batch_height += est_height

        # 마지막 배치 및 Main Disclaimer 처리
        if (layout.top + current_batch_height + DISCLAIMER_MAIN_HEIGHT) <= BODY_BOTTOM_LIMIT:
            self._render_box(layout, current_batch, current_batch_height + DISCLAIMER_MAIN_HEIGHT, BOX_WIDTH, highlight_keywords, main_disclaimer=disclaimer_main)
            self._draw_footer_info(layout, FOOTER_TOP)
        else:
            if current_batch:
                self._render_box(layout, current_batch, current_batch_height, BOX_WIDTH, highlight_keywords, main_disclaimer=None)
                self._draw_footer_info(layout, FOOTER_TOP)
            
            layout.add_new_slide()
            current_page += 1
            self._draw_comment_header(layout, page_info=(current_page, total_pages))
            self._render_box(layout, [], DISCLAIMER_MAIN_HEIGHT, BOX_WIDTH, highlight_keywords, main_disclaimer=disclaimer_main)
            self._draw_footer_info(layout, FOOTER_TOP)

    def _calculate_comment_pages(self, layout, comments, body_limit, line_height, chars_per_line, disclaimer_height):
        """Comment 섹션이 몇 페이지에 걸칠지 미리 계산합니다."""
        if not comments:
            return 1
        
        page_count = 1
        current_top = layout.top
        current_batch_height = 0
        
        for comment in comments:
            text_len = len(comment)
            lines = int((text_len / chars_per_line) + 1)
            est_height = (lines * line_height) + Cm(0.2)
            
            if (current_top + current_batch_height + est_height) > body_limit:
                page_count += 1
                current_top = self.config.BODY_TOP_START
                current_batch_height = Cm(1.0)  # header height
            
            current_batch_height += est_height
        
        # 마지막 페이지에 disclaimer가 들어가는지 확인
        if (current_top + current_batch_height + disclaimer_height) > body_limit:
            page_count += 1
        
        return page_count


    def _render_box(self, layout, comments_batch, height, width, highlight_keywords, main_disclaimer=None):
        if not comments_batch and not main_disclaimer:
            return

        from pptx.dml.color import RGBColor
        from pptx.util import Pt
        
        # 페이지 중앙 정렬
        slide_width = layout.prs.slide_width
        left_position = (slide_width - width) / 2
        
        tb = layout.current_slide.shapes.add_textbox(left_position, layout.top, width, height)
        
        # 테두리 설정 (검은색 0.75pt)
        line = tb.line
        line.color.rgb = RGBColor(0, 0, 0)
        line.width = Pt(0.75)
        
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_top = Cm(0.2)
        tf.margin_bottom = Cm(0.2)
        tf.margin_left = Cm(0.2)
        tf.margin_right = Cm(0.2)
        
        import re
        is_first_paragraph = True
        
        # 코멘트 그리기
        for comment in comments_batch:
            if is_first_paragraph:
                p = tf.paragraphs[0]
                is_first_paragraph = False
            else:
                p = tf.add_paragraph()
                p.space_before = Pt(6) 
            
            p.alignment = PP_ALIGN.JUSTIFY # [Changed] 양쪽 정렬
            p.line_spacing = 1.5 # [Changed] 줄 간격 1.5
            
            
            matched = False
            if highlight_keywords:
                keywords_sorted = sorted(highlight_keywords, key=len, reverse=True)
                for kw in keywords_sorted:
                    if kw in comment:
                         parts = re.split(f"({re.escape(kw)})", comment)
                         for part in parts:
                             is_bold_part = (part == kw)
                             self._set_run_style(p.add_run(), part, is_bold=is_bold_part, 
                                                 font_size=Pt(8), color=self.config.COLOR_BLACK)
                         matched = True
                         break
            
            if not matched:
                match_col = re.match(r"^([^:]+)(:)(.*)$", comment, re.DOTALL)
                if match_col:
                    subject = match_col.group(1)
                    colon = match_col.group(2)
                    rest = match_col.group(3)
                    self._set_run_style(p.add_run(), subject, is_bold=True, 
                                        font_size=Pt(8), color=self.config.COLOR_BLACK)
                    self._set_run_style(p.add_run(), colon + rest, is_bold=False, 
                                        font_size=Pt(8), color=self.config.COLOR_BLACK)
                else:
                    self._set_run_style(p.add_run(), comment, is_bold=False, 
                                        font_size=Pt(8), color=self.config.COLOR_BLACK)
        
        # 메인 고지문 추가
        if main_disclaimer:
            if not is_first_paragraph: # 코멘트가 있었다면 띄움
                p_spacer = tf.add_paragraph()
                p_spacer.space_before = Pt(12)
            
            p_disc = tf.add_paragraph() if not is_first_paragraph else tf.paragraphs[0]
            if not is_first_paragraph:
                p_disc.space_before = Pt(6)
            
            self._set_run_style(p_disc.add_run(), main_disclaimer, font_size=Pt(8), color=self.config.COLOR_BLACK)

        layout.add_space(height)

    def _draw_footer_info(self, layout, top_position):
        """페이지 하단에 고정된 Raw Data 안내 문구 (투명 박스)"""
        from pptx.util import Pt
        
        footer_text_1 = "• 본 검사의 raw data (BAM, FASTQ, VCF) 파일은 분자 병리 검사실 내 병리과 서버 컴퓨터에서 보관, 관리되고 있습니다."
        footer_text_2 = "• 본 검사의 결과는 검체에 포함된 정상세포와 암세포의 비율에 따라 위음성의 결과를 배제 할 수 없습니다."
        
        width = Cm(17.55) 
        height = Cm(1.5)
        
        # [CENTERING] 페이지 중앙 정렬
        slide_width = layout.prs.slide_width
        left_position = (slide_width - width) / 2
        
        tb = layout.current_slide.shapes.add_textbox(left_position, top_position, width, height)
        # 투명 박스 (기본값)
        
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        self._set_run_style(p1.add_run(), footer_text_1, font_size=Pt(8), color=self.config.COLOR_BLACK)
        
        p2 = tf.add_paragraph()
        p2.space_before = Pt(6)
        self._set_run_style(p2.add_run(), footer_text_2, font_size=Pt(8), color=self.config.COLOR_BLACK)

    def _draw_long_text(self, layout, text, font_size=Pt(8), is_bold_keyword=None, highlight_keywords=None):
        height = Cm(0.8)
        layout.check_space(height)
        tb = layout.current_slide.shapes.add_textbox(layout.margin, layout.top, layout.width, height)
        p = tb.text_frame.paragraphs[0]
        
        matched = False
        if highlight_keywords:
            import re
            keywords_sorted = sorted(highlight_keywords, key=len, reverse=True)
            for kw in keywords_sorted:
                if kw in text:
                     parts = re.split(f"({re.escape(kw)})", text)
                     for part in parts:
                         is_bold = (part == kw)
                         # 폰트 사이즈는 파라미터를 따라감
                         self._set_run_style(p.add_run(), part, font_size=font_size, color=self.config.COLOR_BLACK, is_bold=is_bold)
                     matched = True
                     break
        
        if not matched:
            if is_bold_keyword and is_bold_keyword in text:
                 # (호환성을 위해 남겨둠)
                 import re
                 pattern = f"({re.escape(is_bold_keyword)})"
                 parts = re.split(pattern, text)
                 for part in parts:
                     if not part: continue
                     is_bold = (part == is_bold_keyword)
                     self._set_run_style(p.add_run(), part, font_size=font_size, color=self.config.COLOR_BLACK, is_bold=is_bold)
            else:
                 # 2순위: Colon 패턴
                 regex_col = r"^([^:]+)(:)(.*)$"
                 match_col = re.match(regex_col, text, re.DOTALL)
                 
                 if match_col:
                    subject = match_col.group(1)
                    colon = match_col.group(2)
                    rest = match_col.group(3)
                    self._set_run_style(p.add_run(), subject, is_bold=True, font_size=font_size, color=self.config.COLOR_BLACK)
                    self._set_run_style(p.add_run(), colon + rest, is_bold=False, font_size=font_size, color=self.config.COLOR_BLACK)
                 else:
                    self._set_run_style(p.add_run(), text, font_size=font_size, color=self.config.COLOR_BLACK)

        layout.add_space(height)

    def _fill_footer_info(self, slide, report_data):
        """
        마지막 페이지의 템플릿 요소(TextBox, Table)를 찾아 Footer 정보를 기입합니다.
        - 분자 접수 번호: 테이블 (Row 0, Col 1)
        - Tested by, Analyzed by, Signed by: 텍스트 박스 (기존 텍스트 뒤에 Append)
        """
        diagnosis_user = report_data.get('diagnosis_user', {})
        if not diagnosis_user:
            return

        # 데이터 매핑
        receipt_no = diagnosis_user.get('분자접수번호', '')
        user_info_map = {
            "Tested by:": diagnosis_user.get('Tested by', ''),
            "Analyzed by:": diagnosis_user.get('Analyzed by', ''),
            "Signed by:": diagnosis_user.get('Signed by', '')
        }

        for shape in slide.shapes:
            # 1. 텍스트 박스 처리 (Tested by, Analyzed by, Signed by)
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                for key, value in user_info_map.items():
                    # 키워드가 포함되어 있고, 아직 값이 채워지지 않은 경우(길이 체크 등)
                    if key in text and value:
                        # 이미 값이 들어있는지 확인 (중복 방지)
                        if value in text: 
                            continue

                        # 스타일(폰트, 사이즈)을 유지하며 추가하기 위해 Run 추가
                        p = shape.text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = f" {value}"
                        run.font.name = self.config.FONT_NAME
                        # 기존 텍스트의 폰트 크기를 따라가거나, 기본 Body 사이즈 적용
                        # 여기서는 Body 사이즈 적용
                        from pptx.util import Pt
                        run.font.size = Pt(9) 
                        run.font.bold = True # 이름 등은 Bold 처리


            elif shape.has_table:
                tbl = shape.table
                try:
                    # (0, 0) 셀이 "분자 접수 번호" 인지 확인
                    if len(tbl.rows) > 0 and len(tbl.rows[0].cells) > 0:
                        header_text = tbl.rows[0].cells[0].text_frame.text.strip()
                        if "분자 접수 번호" in header_text or "분자접수번호" in header_text:
                            if len(tbl.rows[0].cells) > 1:
                                target_cell = tbl.rows[0].cells[1]
                                self._set_cell_text_preserving_style(
                                    target_cell, 
                                    str(receipt_no), 
                                    is_bold=True, # [Changed] Bold 처리 (사용자 요청)
                                    font_color=self.config.COLOR_BLACK
                                )
                except Exception as e:
                    print(f"Footer table update error: {e}")

    def _fill_diagnostic_info(self, prs, report_data):
        """
        '검사기기' 정보를 찾아 기입합니다.
        검사정보 섹션의 텍스트 박스(예: '검사기기 :')를 찾아 뒤에 값을 이어 붙입니다.
        """
        diag_info = report_data.get('diagnostic_info', {})
        target_value = diag_info.get('검사기기', '')
        
        if not target_value:
            return

        # 모든 슬라이드에서 검색 (위치가 명확하지 않으므로 안전하게 전체 검색)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        # "검사기기" 키워드 포함 여부 확인
                        if "검사기기" in p.text:
                            # 이미 값이 들어있는지 확인 (중복 기입 방지)
                            if target_value in p.text:
                                continue
                            
                            # 값 추가 (공백 한 칸 + 값)
                            run = p.add_run()
                            run.text = f" {target_value}"
                            run.font.name = self.config.FONT_NAME
                            # 기존 텍스트 스타일을 따라가거나 명시적 설정
                            # 사용자가 10pt를 요청함 (기본 Body는 9pt이나 이 섹션은 10pt가 적절해 보임)
                            from pptx.util import Pt
                            run.font.size = Pt(10)
                            run.font.bold = False # 값 부분은 Bold 아님 (HTML 참조: span만 label 클래스)
                            
                            # 하나 찾으면 종료 (일반적으로 하나만 존재)
                            return
